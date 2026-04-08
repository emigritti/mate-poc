"""
Integration Agent — Document Parser
Extracts text content from uploaded best-practice documents for Knowledge Base ingestion.

Supported formats:
  - PDF  (via PyMuPDF / fitz)
  - DOCX (via python-docx)
  - XLSX (via openpyxl)
  - PPTX (via python-pptx)
  - MD   (plain text read)
  - PNG / JPG (via vision_service.caption_figure — LLaVA)
  - SVG  (text extraction via stdlib xml.etree.ElementTree)

Each parser returns the full text content as a string.  The chunker then
splits the text into overlapping fragments ready for ChromaDB insertion.
"""

import io
import logging
import re
from dataclasses import dataclass

logger = logging.getLogger(__name__)

# ── MIME type mapping ─────────────────────────────────────────────────────────
ALLOWED_KB_MIME: dict[str, str] = {
    # PDF
    "application/pdf": "pdf",
    # Word
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/msword": "docx",
    # Excel
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.ms-excel": "xlsx",
    # PowerPoint
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-powerpoint": "pptx",
    # Markdown / plain text
    "text/markdown": "md",
    "text/plain": "md",
    "text/x-markdown": "md",
    # Images
    "image/png":     "png",
    "image/jpeg":    "jpg",
    "image/jpg":     "jpg",
    "image/svg+xml": "svg",
}

# Also accept by file extension (fallback when MIME is unreliable)
ALLOWED_KB_EXTENSIONS: dict[str, str] = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "docx",
    ".xlsx": "xlsx",
    ".xls": "xlsx",
    ".pptx": "pptx",
    ".ppt": "pptx",
    ".md": "md",
    ".txt": "md",
    # Images
    ".png":  "png",
    ".jpg":  "jpg",
    ".jpeg": "jpg",
    ".svg":  "svg",
}

# Image file types handled by the standalone image parser (bypass Docling)
_IMAGE_TYPES: frozenset[str] = frozenset({"png", "jpg", "svg"})


# ── Exceptions ────────────────────────────────────────────────────────────────

class DocumentParseError(ValueError):
    """Raised when a document cannot be parsed."""


# ── Data classes ──────────────────────────────────────────────────────────────

@dataclass
class ParseResult:
    """Result of document parsing."""
    text: str
    file_type: str
    page_count: int  # pages/slides/sheets depending on format


@dataclass
class TextChunk:
    """A single chunk of text ready for vector insertion."""
    text: str
    index: int       # 0-based chunk position
    metadata: dict   # {source_page, ...}


@dataclass
class DoclingChunk:
    """A single chunk produced by the Docling layout-aware parser (ADR-031).

    chunk_type values: "text", "table", "figure"
    section_header:    heading text from the nearest parent section in the document.
    """
    text: str
    chunk_type: str    # "text" | "table" | "figure"
    page_num: int
    section_header: str
    index: int
    metadata: dict


# ── Format-specific parsers ───────────────────────────────────────────────────

def _parse_pdf(data: bytes) -> ParseResult:
    """Extract text from PDF using PyMuPDF."""
    import fitz  # PyMuPDF

    try:
        doc = fitz.open(stream=data, filetype="pdf")
    except Exception as exc:
        raise DocumentParseError(f"Failed to open PDF: {exc}") from exc

    pages: list[str] = []
    for page in doc:
        text = page.get_text("text")
        if text.strip():
            pages.append(text.strip())

    doc.close()

    if not pages:
        raise DocumentParseError("PDF contains no extractable text.")

    return ParseResult(
        text="\n\n".join(pages),
        file_type="pdf",
        page_count=len(pages),
    )


def _parse_docx(data: bytes) -> ParseResult:
    """Extract text from DOCX using python-docx."""
    from docx import Document as DocxDocument

    try:
        doc = DocxDocument(io.BytesIO(data))
    except Exception as exc:
        raise DocumentParseError(f"Failed to open DOCX: {exc}") from exc

    paragraphs: list[str] = []

    # Extract paragraphs
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)

    # Extract tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))

    if not paragraphs:
        raise DocumentParseError("DOCX contains no extractable text.")

    return ParseResult(
        text="\n\n".join(paragraphs),
        file_type="docx",
        page_count=1,  # DOCX doesn't have a simple page concept
    )


def _parse_xlsx(data: bytes) -> ParseResult:
    """Extract text from XLSX using openpyxl."""
    from openpyxl import load_workbook

    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:
        raise DocumentParseError(f"Failed to open XLSX: {exc}") from exc

    sheets: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"## Sheet: {sheet_name}\n" + "\n".join(rows))

    wb.close()

    if not sheets:
        raise DocumentParseError("XLSX contains no extractable text.")

    return ParseResult(
        text="\n\n".join(sheets),
        file_type="xlsx",
        page_count=len(sheets),
    )


def _parse_pptx(data: bytes) -> ParseResult:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise DocumentParseError(f"Failed to open PPTX: {exc}") from exc

    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        # Also extract notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            for para in slide.notes_slide.notes_text_frame.paragraphs:
                note = para.text.strip()
                if note:
                    texts.append(f"[Note] {note}")
        if texts:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))

    if not slides:
        raise DocumentParseError("PPTX contains no extractable text.")

    return ParseResult(
        text="\n\n".join(slides),
        file_type="pptx",
        page_count=len(slides),
    )


def _parse_markdown(data: bytes) -> ParseResult:
    """Read markdown/plain text as-is."""
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = data.decode("latin-1")
        except UnicodeDecodeError as exc:
            raise DocumentParseError(f"Cannot decode text file: {exc}") from exc

    if not text.strip():
        raise DocumentParseError("Markdown file is empty.")

    return ParseResult(
        text=text.strip(),
        file_type="md",
        page_count=1,
    )


# ── Parser dispatch ───────────────────────────────────────────────────────────

_PARSERS: dict[str, callable] = {
    "pdf": _parse_pdf,
    "docx": _parse_docx,
    "xlsx": _parse_xlsx,
    "pptx": _parse_pptx,
    "md": _parse_markdown,
}


def detect_file_type(filename: str, content_type: str | None) -> str:
    """Determine file type from MIME type or extension.

    Returns one of: pdf, docx, xlsx, pptx, md, png, jpg, svg.
    Raises DocumentParseError if not supported.
    """
    # Try MIME type first
    if content_type and content_type in ALLOWED_KB_MIME:
        return ALLOWED_KB_MIME[content_type]

    # Fallback to extension
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext in ALLOWED_KB_EXTENSIONS:
        return ALLOWED_KB_EXTENSIONS[ext]

    raise DocumentParseError(
        f"Unsupported file type: MIME='{content_type}', filename='{filename}'. "
        f"Supported: PDF, DOCX, XLSX, PPTX, MD, PNG, JPG, SVG."
    )


def parse_document(data: bytes, filename: str, content_type: str | None) -> ParseResult:
    """Extract text from a document file.

    Args:
        data:         Raw file bytes.
        filename:     Original filename (used for extension fallback).
        content_type: MIME type from the upload (may be None or unreliable).

    Returns:
        ParseResult with extracted text, detected file type, and page count.

    Raises:
        DocumentParseError: if the file cannot be parsed or is empty.
    """
    file_type = detect_file_type(filename, content_type)
    parser = _PARSERS[file_type]

    logger.info("[KB] Parsing %s as %s (%d bytes)...", filename, file_type, len(data))

    result = parser(data)

    logger.info(
        "[KB] Parsed %s: %d chars, %d pages/sections.",
        filename, len(result.text), result.page_count,
    )
    return result


# ── Chunker ───────────────────────────────────────────────────────────────────

def chunk_text(
    text: str,
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
) -> list[TextChunk]:
    """Split text into overlapping chunks for vector store insertion.

    Uses sentence-aware splitting: tries to break at sentence boundaries
    (period + space) to avoid cutting mid-sentence.

    Args:
        text:          Full document text.
        chunk_size:    Target characters per chunk.
        chunk_overlap: Characters of overlap between consecutive chunks.

    Returns:
        List of TextChunk objects with text and positional metadata.
    """
    if not text.strip():
        return []

    # Normalise whitespace
    clean = re.sub(r"\n{3,}", "\n\n", text.strip())

    chunks: list[TextChunk] = []
    start = 0

    while start < len(clean):
        end = start + chunk_size

        # If not at the end, try to break at a sentence boundary
        if end < len(clean):
            # Look for last sentence-ending punctuation in the chunk
            for sep in [". ", ".\n", "\n\n", "\n", " "]:
                last_sep = clean.rfind(sep, start + chunk_size // 2, end)
                if last_sep != -1:
                    end = last_sep + len(sep)
                    break

        chunk_text_str = clean[start:end].strip()
        if chunk_text_str:
            chunks.append(TextChunk(
                text=chunk_text_str,
                index=len(chunks),
                metadata={"char_start": start, "char_end": end},
            ))

        # Advance by chunk_size - overlap
        start = max(start + 1, end - chunk_overlap)

    logger.info("[KB] Chunked text into %d chunks (size=%d, overlap=%d).", len(chunks), chunk_size, chunk_overlap)
    return chunks


def semantic_chunk(
    text: str,
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
) -> list[TextChunk]:
    """Split text into overlapping chunks respecting semantic boundaries (R11).

    Uses LangChain RecursiveCharacterTextSplitter with separator priority:
      H2 heading → H3 heading → paragraph → newline → sentence → word

    This replaces fixed-size splitting in chunk_text() for new KB uploads.
    chunk_text() is preserved for backward compatibility.

    ADR-030: Semantic chunking with LangChain RecursiveCharacterTextSplitter.
    """
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    if not text.strip():
        return []

    clean = re.sub(r"\n{3,}", "\n\n", text.strip())

    splitter = RecursiveCharacterTextSplitter(
        separators=["\n## ", "\n### ", "\n\n", "\n", ". ", " "],
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
    )

    lc_chunks = splitter.create_documents([clean])

    result: list[TextChunk] = []
    for i, doc in enumerate(lc_chunks):
        stripped = doc.page_content.strip()
        if stripped:
            result.append(TextChunk(
                text=stripped,
                index=len(result),
                metadata={"char_start": 0, "char_end": len(stripped)},
            ))

    logger.info(
        "[KB] Semantic chunked into %d chunks (size=%d, overlap=%d).",
        len(result), chunk_size, chunk_overlap,
    )
    return result


# ── Docling layout-aware parser (ADR-031) ─────────────────────────────────────

def _get_page_num(item) -> int:
    """Extract page number from a Docling item's provenance list."""
    try:
        return item.prov[0].page_no
    except (AttributeError, IndexError):
        return 0


def _is_section_header(item) -> bool:
    """Return True if the item is a section/heading item."""
    try:
        return item.label.value == "section_header"
    except AttributeError:
        return False


def _is_table_item(item) -> bool:
    """Return True if the item is a Docling TableItem."""
    return type(item).__name__ == "TableItem"


def _is_picture_item(item) -> bool:
    """Return True if the item is a Docling PictureItem."""
    return type(item).__name__ == "PictureItem"


def _pil_to_bytes(pil_image) -> bytes:
    """Convert a PIL image to PNG bytes."""
    import io as _io
    buf = _io.BytesIO()
    pil_image.save(buf, format="PNG")
    return buf.getvalue()


def _extract_svg_text(svg_bytes: bytes) -> str:
    """Extract all visible text from an SVG file using stdlib XML parsing.

    Iterates every element in the SVG tree and collects non-empty .text and
    .tail strings (covers <title>, <desc>, <text>, <tspan>, and similar).
    Returns a placeholder when the file contains no text or is malformed XML.
    """
    import xml.etree.ElementTree as ET
    try:
        root = ET.fromstring(svg_bytes)
        parts: list[str] = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                parts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                parts.append(elem.tail.strip())
        return " ".join(parts) if parts else "[SVG: no text content]"
    except ET.ParseError as exc:
        logger.warning("[KB] SVG XML parse error: %s", exc)
        return "[SVG: parse error]"


async def _parse_image_standalone(file_bytes: bytes, file_type: str) -> list[DoclingChunk]:
    """Parse a standalone image file into a single DoclingChunk.

    PNG/JPG: captions the image via vision_service.caption_figure() (LLaVA).
             Falls back to a placeholder when vision is disabled or fails.
    SVG:     extracts text nodes from the XML tree (stdlib — no extra deps).
    """
    if file_type == "svg":
        text = _extract_svg_text(file_bytes)
        chunk_type = "text"
    else:
        from services.vision_service import caption_figure
        text = await caption_figure(file_bytes)
        chunk_type = "figure"

    return [DoclingChunk(
        text=text,
        chunk_type=chunk_type,
        page_num=1,
        section_header="",
        index=0,
        metadata={},
    )]


async def parse_with_docling(file_bytes: bytes, file_type: str) -> list[DoclingChunk]:
    """Parse a document using IBM Docling for layout-aware chunking (ADR-031).

    Extracts text, table, and figure chunks with rich metadata:
      - chunk_type: "text" | "table" | "figure"
      - page_num: source page number
      - section_header: nearest parent section heading
      - index: 0-based sequential index across all chunks

    Standalone image files (png, jpg, svg) bypass Docling entirely and are
    handled by _parse_image_standalone().

    Falls back to legacy text-only parsing when Docling is not installed.
    Figure captions are generated via vision_service.caption_figure() (llava:7b).
    """
    # Standalone images bypass Docling — handle them directly.
    if file_type in _IMAGE_TYPES:
        return await _parse_image_standalone(file_bytes, file_type)

    import asyncio as _asyncio
    import io as _io

    # Lazy import — allows graceful fallback if Docling is not installed.
    try:
        from docling.document_converter import DocumentConverter
        from docling.datamodel.base_models import DocumentStream
    except ImportError:
        logger.warning(
            "[Docling] Not available — falling back to text-only parser. "
            "Install with: pip install docling"
        )
        return _docling_fallback(file_bytes, file_type)

    # Import vision_service here to avoid circular imports at module level.
    from services.vision_service import caption_figure

    def _convert():
        converter = DocumentConverter()
        stream = DocumentStream(
            name=f"document.{file_type}",
            stream=_io.BytesIO(file_bytes),
        )
        return converter.convert(stream)

    # Docling conversion is CPU-bound — run in thread pool to avoid blocking.
    # A configurable timeout prevents 504s on very large documents (e.g. 500-page books).
    # On timeout, fall back to the fast legacy text parser rather than failing the upload.
    from config import settings as _settings
    loop = _asyncio.get_event_loop()
    try:
        result = await _asyncio.wait_for(
            loop.run_in_executor(None, _convert),
            timeout=_settings.docling_timeout_seconds,
        )
    except _asyncio.TimeoutError:
        logger.warning(
            "[Docling] Parsing timed out after %ds — falling back to legacy parser. "
            "Set DOCLING_TIMEOUT_SECONDS env var to allow more time for large documents.",
            _settings.docling_timeout_seconds,
        )
        return _docling_fallback(file_bytes, file_type)
    doc = result.document

    chunks: list[DoclingChunk] = []
    current_section = ""
    chunk_idx = 0

    for item, _level in doc.iterate_items():
        page_num = _get_page_num(item)

        if _is_section_header(item):
            current_section = getattr(item, "text", "")
            continue

        if _is_table_item(item):
            table_md = item.export_to_markdown() or ""
            if table_md.strip():
                chunks.append(DoclingChunk(
                    text=f"[TABLE]\n{table_md.strip()}",
                    chunk_type="table",
                    page_num=page_num,
                    section_header=current_section,
                    index=chunk_idx,
                    metadata={},
                ))
                chunk_idx += 1

        elif _is_picture_item(item):
            try:
                pil_img = item.get_image(doc)
                img_bytes = _pil_to_bytes(pil_img) if pil_img else None
            except Exception:
                img_bytes = None

            caption = await caption_figure(img_bytes) if img_bytes is not None else "[FIGURE: no caption available]"
            chunks.append(DoclingChunk(
                text=caption,
                chunk_type="figure",
                page_num=page_num,
                section_header=current_section,
                index=chunk_idx,
                metadata={},
            ))
            chunk_idx += 1

        else:
            # Default: treat as text item
            text = getattr(item, "text", "").strip()
            if text:
                chunks.append(DoclingChunk(
                    text=text,
                    chunk_type="text",
                    page_num=page_num,
                    section_header=current_section,
                    index=chunk_idx,
                    metadata={},
                ))
                chunk_idx += 1

    logger.info(
        "[Docling] Parsed %d chunks (%d text, %d table, %d figure) from %s.",
        len(chunks),
        sum(1 for c in chunks if c.chunk_type == "text"),
        sum(1 for c in chunks if c.chunk_type == "table"),
        sum(1 for c in chunks if c.chunk_type == "figure"),
        file_type,
    )
    return chunks


def _docling_fallback(file_bytes: bytes, file_type: str) -> list[DoclingChunk]:
    """Wrap legacy text parsing into DoclingChunk list when Docling is unavailable."""
    if file_type == "md":
        try:
            text = file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            text = file_bytes.decode("latin-1", errors="replace")
        if not text.strip():
            return []
        text_chunks = semantic_chunk(text)
    else:
        try:
            parse_result = _PARSERS[file_type](file_bytes)
            text_chunks = semantic_chunk(parse_result.text)
        except Exception as exc:
            logger.warning("[Docling-fallback] Parse failed: %s", exc)
            return []

    return [
        DoclingChunk(
            text=tc.text,
            chunk_type="text",
            page_num=0,
            section_header="",
            index=i,
            metadata=tc.metadata,
        )
        for i, tc in enumerate(text_chunks)
    ]
